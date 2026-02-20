"""Document parser — PDF, DOCX, image, text → structured markdown."""

from __future__ import annotations

from pathlib import Path

from app.models.internal import ParsedDocument
from app.utils.exceptions import DocumentParsingError
from app.utils.logger import get_logger

log = get_logger("document_parser")


class DocumentParser:
    """Converts PDF/DOCX/image/text into clean structured markdown."""

    async def parse(self, file_path: Path, mime_type: str) -> ParsedDocument:
        """Route to correct parser based on MIME type."""
        try:
            if mime_type == "application/pdf":
                return await self.parse_pdf(file_path)
            elif mime_type in (
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ):
                return await self.parse_docx(file_path)
            elif mime_type == "application/msword":
                return await self.parse_doc(file_path)
            elif mime_type in (
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ):
                return await self.parse_pptx(file_path)
            elif mime_type.startswith("image/"):
                return await self.parse_image_ocr(file_path)
            elif mime_type == "text/plain":
                return await self.parse_text(file_path)
            else:
                raise DocumentParsingError(f"Unsupported MIME type: {mime_type}")
        except DocumentParsingError:
            raise
        except Exception as e:
            raise DocumentParsingError(f"Parsing failed: {e}") from e

    async def parse_pdf(self, path: Path) -> ParsedDocument:
        """Extract text from PDF using pymupdf4llm, fallback to OCR."""
        import pymupdf4llm

        md_text = pymupdf4llm.to_markdown(str(path))

        if len(md_text.strip()) < 100:
            log.info("pdf_low_text_fallback_ocr", path=str(path))
            return await self.parse_image_ocr(path)

        page_count = 1
        try:
            import pymupdf
            doc = pymupdf.open(str(path))
            page_count = len(doc)
            doc.close()
        except Exception:
            pass

        lang = self._detect_language(md_text)
        log.info("pdf_parsed", pages=page_count, chars=len(md_text), lang=lang)

        return ParsedDocument(
            text=md_text,
            markdown=md_text,
            page_count=page_count,
            language=lang,
            mime_type="application/pdf",
        )

    async def parse_docx(self, path: Path) -> ParsedDocument:
        """Extract text from DOCX preserving structure."""
        from docx import Document

        doc = Document(str(path))
        paragraphs = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style = para.style.name.lower() if para.style else ""
            if "heading" in style:
                level = 1
                for c in style:
                    if c.isdigit():
                        level = int(c)
                        break
                paragraphs.append(f"{'#' * level} {text}")
            else:
                paragraphs.append(text)

        # Extract tables
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append("| " + " | ".join(cells) + " |")
            if rows:
                header_sep = "| " + " | ".join(["---"] * len(table.rows[0].cells)) + " |"
                rows.insert(1, header_sep)
                paragraphs.append("\n".join(rows))

        md_text = "\n\n".join(paragraphs)
        lang = self._detect_language(md_text)

        log.info("docx_parsed", paras=len(paragraphs), chars=len(md_text), lang=lang)
        return ParsedDocument(
            text=md_text,
            markdown=md_text,
            page_count=max(1, len(paragraphs) // 30),
            language=lang,
            mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )

    async def parse_doc(self, path: Path) -> ParsedDocument:
        """Extract text from legacy .doc files."""
        import subprocess

        try:
            # Try antiword first (most reliable for .doc)
            result = subprocess.run(
                ["antiword", str(path)],
                capture_output=True, text=True, timeout=30
            )
            if result.returncode == 0 and result.stdout.strip():
                text = result.stdout
            else:
                # Fallback: catdoc
                result = subprocess.run(
                    ["catdoc", str(path)],
                    capture_output=True, text=True, timeout=30
                )
                if result.returncode == 0:
                    text = result.stdout
                else:
                    raise DocumentParsingError("Cannot parse .doc file. Install 'antiword': brew install antiword")

            lang = self._detect_language(text)
            log.info("doc_parsed", chars=len(text), lang=lang)
            return ParsedDocument(
                text=text, markdown=text, page_count=max(1, len(text) // 3000),
                language=lang, mime_type="application/msword",
            )
        except subprocess.TimeoutExpired:
            raise DocumentParsingError(".doc parsing timed out")
        except FileNotFoundError:
            raise DocumentParsingError(
                "Cannot parse .doc files — install antiword: brew install antiword"
            )

    async def parse_pptx(self, path: Path) -> ParsedDocument:
        """Extract text from PPTX slides."""
        from pptx import Presentation

        prs = Presentation(str(path))
        slides_text = []

        for i, slide in enumerate(prs.slides, 1):
            slide_parts = [f"## Slide {i}"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_parts.append(text)
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append("| " + " | ".join(cells) + " |")
                    if rows:
                        header_sep = "| " + " | ".join(["---"] * len(table.rows[0].cells)) + " |"
                        rows.insert(1, header_sep)
                        slide_parts.append("\n".join(rows))
            slides_text.append("\n".join(slide_parts))

        md_text = "\n\n".join(slides_text)
        lang = self._detect_language(md_text)

        log.info("pptx_parsed", slides=len(prs.slides), chars=len(md_text), lang=lang)
        return ParsedDocument(
            text=md_text, markdown=md_text, page_count=len(prs.slides),
            language=lang,
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

    async def parse_image_ocr(self, path: Path) -> ParsedDocument:
        """OCR for scanned documents using EasyOCR."""
        try:
            import easyocr
            reader = easyocr.Reader(["en", "hi", "mr", "ta"], gpu=False)
            results = reader.readtext(str(path))
            text = "\n".join([r[1] for r in results])

            lang = self._detect_language(text)
            log.info("ocr_parsed", chars=len(text), lang=lang)

            return ParsedDocument(
                text=text,
                markdown=text,
                page_count=1,
                language=lang,
                mime_type="image/ocr",
            )
        except Exception as e:
            raise DocumentParsingError(f"OCR failed: {e}") from e

    async def parse_text(self, path: Path) -> ParsedDocument:
        """Parse plain text file."""
        text = path.read_text(encoding="utf-8", errors="replace")
        lang = self._detect_language(text)
        return ParsedDocument(
            text=text, markdown=text, page_count=1, language=lang, mime_type="text/plain"
        )

    def _detect_language(self, text: str) -> str:
        """Detect language of document text."""
        try:
            from langdetect import detect
            return detect(text[:2000])
        except Exception:
            return "en"

    def detect_contract_type(self, text: str) -> str:
        """Keyword-based contract type classification."""
        text_lower = text.lower()
        patterns = {
            "rental": ["tenant", "landlord", "rent", "lease", "tenancy", "premises", "security deposit"],
            "employment": ["employee", "employer", "salary", "designation", "probation", "termination of employment"],
            "freelance": ["client", "deliverable", "invoice", "scope of work", "independent contractor"],
            "loan": ["borrower", "lender", "emi", "interest rate", "principal amount", "disbursement"],
            "nda": ["confidential", "non-disclosure", "proprietary information", "trade secret"],
            "startup": ["founder", "equity", "vesting", "investor", "valuation", "cap table"],
            "consumer": ["consumer", "warranty", "refund", "product liability", "defect"],
        }

        scores = {}
        for ctype, keywords in patterns.items():
            score = sum(1 for kw in keywords if kw in text_lower)
            if score > 0:
                scores[ctype] = score

        if scores:
            return max(scores, key=scores.get)
        return "general"
